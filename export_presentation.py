"""
export_presentation.py
Converts the Zentai Networks slide deck (backend/presentation_slides.py)
into a fully formatted PowerPoint file (.pptx) that can be opened in
Microsoft PowerPoint, Keynote, or Google Slides.

Run this from the project root:
    python3 export_presentation.py
"""

import re
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Make sure we can import the slides from the backend ──────────────────────
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from backend.presentation_slides import SLIDES


# ── Brand colours (matching the Cyberpunk UI) ────────────────────────────────
CYAN        = RGBColor(0x00, 0xE5, 0xFF)    # #00E5FF  — accent headings
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)    # #8B5CF6  — secondary accents
DARK_BG     = RGBColor(0x0D, 0x11, 0x17)    # #0D1117  — slide background
TEXT_LIGHT  = RGBColor(0xE6, 0xED, 0xF3)    # #E6EDF3  — body text
TEXT_DIM    = RGBColor(0xCC, 0xCC, 0xCC)    # #CCCCCC  — dimmed text
ACCENT_LINE = RGBColor(0x21, 0x26, 0x2E)    # #21262E  — separator lines

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def strip_html(html: str) -> str:
    """
    Strip HTML tags from slide content and convert common tags into
    clean plain text for PowerPoint rendering.
    Handles: <ul><li>, <ol><li>, <b>, <blockquote>, <table><tr><td>
    """
    # Convert list items into bullet dashes
    html = re.sub(r'<li>', '• ', html)
    html = re.sub(r'<\/?ul>|<\/?ol>', '\n', html)

    # Bold text → wrap in asterisks (we strip them later, just for readability)
    html = re.sub(r'<b>(.*?)</b>', r'\1', html, flags=re.DOTALL)
    html = re.sub(r'<i>(.*?)</i>', r'\1', html, flags=re.DOTALL)

    # Blockquotes → indent with a quote mark
    html = re.sub(r'<blockquote>(.*?)</blockquote>', r'❝ \1', html, flags=re.DOTALL)

    # Table rows → convert each cell to a dash-separated line
    html = re.sub(r'<th[^>]*>(.*?)</th>', r'\1  |  ', html, flags=re.DOTALL)
    html = re.sub(r'<td[^>]*>(.*?)</td>', r'\1  |  ', html, flags=re.DOTALL)
    html = re.sub(r'<tr[^>]*>', '\n', html)
    html = re.sub(r'<\/?tr>|<\/?table[^>]*>|<\/?thead>|<\/?tbody>', '', html)

    # Ordered list numbers
    html = re.sub(r'<li>', '  ', html)

    # Line breaks
    html = re.sub(r'<br\s*/?>', '\n', html)
    html = re.sub(r'<hr[^>]*>', '\n─────────────────────────────────\n', html)

    # Code tags
    html = re.sub(r'<code>(.*?)</code>', r'`\1`', html)

    # Remove all remaining HTML tags
    html = re.sub(r'<[^>]+>', '', html)

    # Collapse multiple blank lines into a single blank line
    html = re.sub(r'\n{3,}', '\n\n', html)

    return html.strip()


def add_slide(prs: Presentation, slide_number: int, title: str, content: str):
    """
    Adds a single formatted slide to the presentation.

    Layout:
        - Full dark background rectangle
        - Slide number badge (top-left corner, small)
        - Title (large cyan heading)
        - Horizontal accent rule
        - Body text (clean, readable)
        - Company name watermark (bottom-right, dimmed)
    """
    blank_layout = prs.slide_layouts[6]   # "Blank" layout — no default placeholders
    slide = prs.slides.add_slide(blank_layout)

    # ── Dark background ───────────────────────────────────────────────────────
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.color.rgb = DARK_BG

    # ── Slide number badge (top-left) ─────────────────────────────────────────
    badge = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(1.5), Inches(0.4))
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"  {slide_number} / {len(SLIDES)}  "
    run.font.size = Pt(10)
    run.font.color.rgb = CYAN
    run.font.bold = True

    # ── Title ─────────────────────────────────────────────────────────────────
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35),
        Inches(12.33), Inches(1.1)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = CYAN

    # ── Horizontal accent divider ─────────────────────────────────────────────
    divider = slide.shapes.add_shape(
        1,
        Inches(0.5), Inches(1.5),
        Inches(12.33), Pt(2)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = CYAN
    divider.line.color.rgb = CYAN

    # ── Body text ─────────────────────────────────────────────────────────────
    clean_body = strip_html(content)

    body_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.65),
        Inches(12.33), Inches(5.4)
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    lines = clean_body.split('\n')
    first = True
    for line in lines:
        line = line.strip()
        if not line:
            # Add vertical whitespace between sections
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            p.space_after = Pt(4)
            first = False
            continue

        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False

        run = p.add_run()
        run.text = line
        run.font.color.rgb = TEXT_LIGHT
        run.font.size = Pt(17)

        # Slightly dim bullet lines
        if line.startswith("•"):
            run.font.size = Pt(16)
            p.space_before = Pt(4)

        # Highlight blockquote lines in italic cyan
        if line.startswith("❝"):
            run.font.color.rgb = CYAN
            run.font.italic = True
            run.font.size = Pt(18)

        # Separator lines stay dim
        if "─────" in line:
            run.font.color.rgb = ACCENT_LINE
            run.font.size = Pt(10)

    # ── Watermark (bottom-right) ──────────────────────────────────────────────
    wm = slide.shapes.add_textbox(
        Inches(9.5), Inches(6.9),
        Inches(3.5), Inches(0.4)
    )
    tf = wm.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Zentai Networks — Capstone 2025"
    run.font.size = Pt(9)
    run.font.color.rgb = TEXT_DIM
    run.font.italic = True


def main():
    output_path = "Zentai_Networks_Capstone_Presentation.pptx"

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print(f"Generating {len(SLIDES)}-slide presentation…")

    for i, slide_data in enumerate(SLIDES, start=1):
        title   = slide_data.get("title", "")
        content = slide_data.get("content", "")
        add_slide(prs, i, title, content)
        print(f"  ✓ Slide {i:02d}: {title}")

    prs.save(output_path)
    print(f"\nDone! Saved to → {output_path}")
    print("Open it with Microsoft PowerPoint, Keynote, or Google Slides.")


if __name__ == "__main__":
    main()
